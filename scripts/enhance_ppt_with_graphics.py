from __future__ import annotations

from math import cos, pi, sin
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE_PPT = ROOT / "Final_Presentation_Polished.pptx"
OUTPUT_PPT = ROOT / "Final_Presentation_Official_Graphics.pptx"
ASSET_DIR = ROOT / "artifacts" / "presentation_graphics"


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def make_gradient(size: tuple[int, int], left: tuple[int, int, int], right: tuple[int, int, int]) -> Image.Image:
    width, height = size
    img = Image.new("RGBA", size, (0, 0, 0, 0))
    px = img.load()
    for x in range(width):
        blend = x / max(width - 1, 1)
        color = tuple(int(left[i] * (1 - blend) + right[i] * blend) for i in range(3)) + (255,)
        for y in range(height):
            px[x, y] = color
    return img


def create_bearing_graphic(path: Path) -> None:
    size = (1400, 1000)
    img = Image.new("RGBA", size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    cx, cy = 700, 500

    shadow = Image.new("RGBA", size, (0, 0, 0, 0))
    sdraw = ImageDraw.Draw(shadow)
    sdraw.ellipse((165, 95, 1235, 905), fill=(0, 0, 0, 80))
    shadow = shadow.filter(ImageFilter.GaussianBlur(28))
    img.alpha_composite(shadow)

    outer_fill = make_gradient(size, (226, 232, 240), (148, 163, 184))
    inner_fill = make_gradient(size, (241, 245, 249), (100, 116, 139))
    accent_fill = make_gradient(size, (250, 204, 21), (245, 158, 11))

    ring_mask = Image.new("L", size, 0)
    mdraw = ImageDraw.Draw(ring_mask)
    mdraw.ellipse((140, 70, 1260, 930), fill=255)
    mdraw.ellipse((290, 220, 1110, 780), fill=0)
    img.paste(outer_fill, (0, 0), ring_mask)
    draw.ellipse((140, 70, 1260, 930), outline=(71, 85, 105, 255), width=18)
    draw.ellipse((290, 220, 1110, 780), outline=(71, 85, 105, 255), width=18)

    core_mask = Image.new("L", size, 0)
    mdraw = ImageDraw.Draw(core_mask)
    mdraw.ellipse((360, 290, 1040, 710), fill=255)
    mdraw.ellipse((500, 430, 900, 570), fill=0)
    img.paste(inner_fill, (0, 0), core_mask)
    draw.ellipse((360, 290, 1040, 710), outline=(51, 65, 85, 255), width=14)
    draw.ellipse((500, 430, 900, 570), outline=(51, 65, 85, 255), width=12)

    channel_mask = Image.new("L", size, 0)
    mdraw = ImageDraw.Draw(channel_mask)
    mdraw.ellipse((255, 185, 1145, 815), fill=255)
    mdraw.ellipse((390, 320, 1010, 680), fill=0)
    channel = Image.new("RGBA", size, (0, 0, 0, 0))
    cdraw = ImageDraw.Draw(channel)
    cdraw.ellipse((255, 185, 1145, 815), fill=(255, 255, 255, 14))
    channel = channel.filter(ImageFilter.GaussianBlur(6))
    img.alpha_composite(channel)

    for i in range(10):
        angle = (2 * pi / 10) * i - pi / 10
        bx = cx + int(305 * cos(angle))
        by = cy + int(205 * sin(angle))
        ball_bounds = (bx - 74, by - 74, bx + 74, by + 74)
        ball_shadow = Image.new("RGBA", size, (0, 0, 0, 0))
        bs = ImageDraw.Draw(ball_shadow)
        bs.ellipse((bx - 76, by - 70, bx + 76, by + 82), fill=(0, 0, 0, 55))
        ball_shadow = ball_shadow.filter(ImageFilter.GaussianBlur(10))
        img.alpha_composite(ball_shadow)
        ball_mask = Image.new("L", size, 0)
        bm = ImageDraw.Draw(ball_mask)
        bm.ellipse(ball_bounds, fill=255)
        img.paste(accent_fill, (0, 0), ball_mask)
        draw.ellipse(ball_bounds, outline=(120, 53, 15, 255), width=10)
        draw.ellipse((bx - 38, by - 50, bx + 12, by - 4), fill=(255, 246, 200, 88))

    draw.rounded_rectangle((930, 95, 1270, 175), radius=26, fill=(15, 23, 42, 214))
    draw.text((972, 115), "Bearing", fill=(255, 255, 255, 255))
    draw.rounded_rectangle((970, 790, 1290, 870), radius=26, fill=(15, 23, 42, 214))
    draw.text((1002, 810), "Rolling Elements", fill=(255, 255, 255, 255))

    img.save(path)


def create_rotor_graphic(path: Path) -> None:
    size = (1600, 800)
    img = Image.new("RGBA", size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)

    shadow = Image.new("RGBA", size, (0, 0, 0, 0))
    sdraw = ImageDraw.Draw(shadow)
    sdraw.rounded_rectangle((110, 225, 1490, 595), radius=90, fill=(0, 0, 0, 55))
    shadow = shadow.filter(ImageFilter.GaussianBlur(20))
    img.alpha_composite(shadow)

    draw.rounded_rectangle((120, 240, 430, 580), radius=54, fill=(30, 41, 59, 245), outline=(148, 163, 184, 255), width=8)
    for x in range(160, 390, 58):
        draw.line((x, 255, x, 565), fill=(71, 85, 105, 220), width=4)

    draw.rounded_rectangle((390, 370, 1310, 445), radius=36, fill=(148, 163, 184, 255), outline=(71, 85, 105, 255), width=8)
    draw.rounded_rectangle((430, 345, 500, 470), radius=28, fill=(226, 232, 240, 255), outline=(51, 65, 85, 255), width=6)
    draw.rounded_rectangle((1130, 345, 1200, 470), radius=28, fill=(226, 232, 240, 255), outline=(51, 65, 85, 255), width=6)

    for base_x in (560, 960):
        draw.ellipse((base_x, 314, base_x + 160, 500), fill=(248, 250, 252, 255), outline=(71, 85, 105, 255), width=6)
        draw.ellipse((base_x + 28, 342, base_x + 132, 472), fill=(15, 23, 42, 245), outline=(100, 116, 139, 255), width=5)
        for i in range(8):
            angle = (2 * pi / 8) * i
            bx = base_x + 80 + int(57 * cos(angle))
            by = 407 + int(57 * sin(angle))
            draw.ellipse((bx - 16, by - 16, bx + 16, by + 16), fill=(245, 158, 11, 255), outline=(120, 53, 15, 255), width=3)

    draw.rounded_rectangle((1255, 285, 1440, 355), radius=20, fill=(14, 165, 233, 235), outline=(8, 47, 73, 255), width=4)
    draw.rectangle((1335, 355, 1360, 430), fill=(8, 47, 73, 255))
    draw.rectangle((1320, 430, 1375, 470), fill=(8, 47, 73, 255))
    points = []
    for i in range(12):
        x = 1200 + i * 24
        y = 270 - (35 if i % 2 else 8)
        points.append((x, y))
    draw.line(points, fill=(16, 185, 129, 255), width=7)
    draw.rounded_rectangle((1215, 165, 1465, 235), radius=18, fill=(15, 23, 42, 220))
    draw.text((1250, 185), "Sensor Output", fill=(255, 255, 255, 255))

    draw.rounded_rectangle((215, 620, 470, 690), radius=18, fill=(15, 23, 42, 220))
    draw.text((245, 640), "Motor Housing", fill=(255, 255, 255, 255))
    draw.rounded_rectangle((592, 620, 890, 690), radius=18, fill=(15, 23, 42, 220))
    draw.text((620, 640), "Bearing + Shaft Zone", fill=(255, 255, 255, 255))

    img.save(path)


def create_fault_badge(path: Path) -> None:
    size = (900, 900)
    img = Image.new("RGBA", size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    draw.ellipse((78, 78, 822, 822), fill=(255, 255, 255, 215), outline=(15, 23, 42, 255), width=18)
    draw.ellipse((180, 180, 720, 720), fill=(241, 245, 249, 255), outline=(100, 116, 139, 255), width=14)
    draw.ellipse((305, 305, 595, 595), fill=(15, 23, 42, 245), outline=(148, 163, 184, 255), width=12)
    for i in range(12):
        angle = (2 * pi / 12) * i
        bx = 450 + int(220 * cos(angle))
        by = 450 + int(220 * sin(angle))
        draw.ellipse((bx - 30, by - 30, bx + 30, by + 30), fill=(249, 115, 22, 255), outline=(124, 45, 18, 255), width=4)
    waveform = []
    start_x = 120
    for i in range(14):
        x = start_x + i * 48
        y = 710 if i % 2 == 0 else 620
        waveform.append((x, y))
    draw.line(waveform, fill=(14, 165, 233, 255), width=14, joint="curve")
    draw.rounded_rectangle((232, 94, 668, 176), radius=24, fill=(15, 23, 42, 230))
    draw.text((268, 117), "Fault Diagnosis", fill=(255, 255, 255, 255))
    img.save(path)


def replace_standalone_text(slide, old: str, new: str) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.text.strip() == old:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = new
            run.font.name = "Times New Roman"
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = RGBColor(15, 23, 42)


def add_caption_tag(slide, left: float, top: float, width: float, text: str) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.33))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    box.fill.transparency = 0.12
    box.line.color.rgb = RGBColor(15, 23, 42)
    box.line.transparency = 0.35
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Times New Roman"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def build_assets() -> dict[str, Path]:
    ensure_dir(ASSET_DIR)
    bearing = ASSET_DIR / "bearing_cross_section.png"
    rotor = ASSET_DIR / "rotating_machine.png"
    badge = ASSET_DIR / "bearing_fault_badge.png"
    create_bearing_graphic(bearing)
    create_rotor_graphic(rotor)
    create_fault_badge(badge)
    return {"bearing": bearing, "rotor": rotor, "badge": badge}


def style_text(prs: Presentation) -> None:
    replacements = {
        2: {"📊": "TIME", "🔍": "FFT"},
        4: {"🏭": "01", "💨": "02", "🚀": "03"},
        6: {"🔒": "01", "🖥️": "02", "⚖️": "03"},
    }
    for slide_idx, mapping in replacements.items():
        slide = prs.slides[slide_idx]
        for old, new in mapping.items():
            replace_standalone_text(slide, old, new)


def add_graphics(prs: Presentation, assets: dict[str, Path]) -> None:
    slide1 = prs.slides[0]
    slide1.shapes.add_picture(str(assets["rotor"]), Inches(8.25), Inches(4.75), width=Inches(4.45))
    slide1.shapes.add_picture(str(assets["bearing"]), Inches(9.85), Inches(0.7), width=Inches(2.35))
    add_caption_tag(slide1, 9.8, 3.1, 2.35, "Rotating Machinery Context")

    slide3 = prs.slides[2]
    slide3.shapes.add_picture(str(assets["bearing"]), Inches(10.05), Inches(3.92), width=Inches(2.72))
    add_caption_tag(slide3, 9.95, 6.55, 2.4, "Bearing Element Model")

    slide9 = prs.slides[8]
    slide9.shapes.add_picture(str(assets["rotor"]), Inches(8.55), Inches(4.8), width=Inches(4.15))
    add_caption_tag(slide9, 8.55, 6.72, 2.45, "Machine-Sensor Pipeline")

    slide15 = prs.slides[14]
    slide15.shapes.add_picture(str(assets["badge"]), Inches(10.75), Inches(1.18), width=Inches(1.65))


def main() -> None:
    assets = build_assets()
    prs = Presentation(str(SOURCE_PPT))
    style_text(prs)
    add_graphics(prs, assets)
    prs.save(str(OUTPUT_PPT))
    print(f"Wrote {OUTPUT_PPT}")
    for name, path in assets.items():
        print(f"{name}: {path}")


if __name__ == "__main__":
    main()
